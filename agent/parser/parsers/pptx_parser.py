# """
# agent/parser/parsers/pptx_parser.py

# PPT 解析器。
# BOO-63 要求：
# - 每页 Slide 输出：标题、bullets、notes
# - 每条要点带 loc（至少 slide 编号）
# - 不依赖大模型主观总结，先输出结构化要点
# - 图片中文字不会被抽取 → 发 warning
# """
# from __future__ import annotations

# import logging
# from pathlib import Path
# from typing import Any

# from agent.parser.base import BaseParser
# from agent.parser.schema import (
#     Error, ImageBlock, Loc, LocType, ParsedDocument, SlideContent,
#     TableBlock, TableMeta, Warning,
# )
# from agent.parser.utils.hash_utils import file_sha1

# logger = logging.getLogger(__name__)


# def _get_placeholder_title(slide) -> str:
#     """
#     安全地从 slide 的占位符中提取标题。
#     python-pptx 在遍历 placeholders 时可能对非占位符形状
#     抛出 ValueError: shape is not a placeholder，所以必须 try/except。
#     """
#     try:
#         for ph in slide.placeholders:
#             if ph.placeholder_format.idx == 0:  # idx 0 = 标题
#                 text = ph.text_frame.text.strip() if ph.has_text_frame else ""
#                 if text:
#                     return text
#     except (ValueError, KeyError, AttributeError):
#         pass
#     return ""


# class PptxParser(BaseParser):
#     supported_extensions = [".pptx"]

#     def __init__(self, kb_name: str = ""):
#         self.kb_name = kb_name

#     def parse(self, file_path: str, **kwargs) -> ParsedDocument:
#         from pptx import Presentation
#         from pptx.util import Inches, Pt
#         from pptx.enum.shapes import MSO_SHAPE_TYPE

#         doc_id = file_sha1(file_path, prefix=self.kb_name)
#         fpath = Path(file_path)
#         warnings: list[Warning] = []
#         errors: list[Error] = []

#         try:
#             prs = Presentation(file_path)
#         except Exception as exc:
#             return ParsedDocument(
#                 doc_id=doc_id,
#                 source=self._source_dict(fpath),
#                 content_type="pptx",
#                 errors=[Error(
#                     error_code="PPTX_OPEN_FAILED",
#                     message=str(exc),
#                     loc=Loc(type=LocType.SLIDE.value, value=0),
#                 )],
#             )

#         slides_data: list[SlideContent] = []
#         all_tables: list[TableBlock] = []
#         all_images: list[ImageBlock] = []
#         title = ""
#         image_count = 0

#         for slide_num, slide in enumerate(prs.slides, start=1):
#             slide_loc = Loc(type=LocType.SLIDE.value, value=slide_num)

#             # ---- 先用安全方法提取占位符标题 ----
#             slide_title = _get_placeholder_title(slide)

#             bullets: list[str] = []
#             text_blocks: list[str] = []
#             slide_tables: list[TableBlock] = []

#             for shape in slide.shapes:
#                 # ---- 文本框 ----
#                 if shape.has_text_frame:
#                     # 判断此 shape 是否为占位符标题（安全检查）
#                     is_title_shape = False
#                     if shape.is_placeholder:
#                         try:
#                             if shape.placeholder_format.idx == 0:
#                                 is_title_shape = True
#                         except (ValueError, AttributeError):
#                             pass

#                     if is_title_shape:
#                         # 标题已经通过 _get_placeholder_title 提取过了
#                         # 如果之前没拿到，这里补上
#                         if not slide_title:
#                             slide_title = shape.text_frame.text.strip()
#                     else:
#                         for para in shape.text_frame.paragraphs:
#                             text = para.text.strip()
#                             if not text:
#                                 continue
#                             # 判断是否为 bullet（有缩进 / 有 bullet 字符）
#                             if para.level > 0 or text.startswith(("•", "-", "·", "▪", "●")):
#                                 bullets.append(text)
#                             else:
#                                 # 如果不是标题文本，加入 text_blocks
#                                 if text != slide_title:
#                                     text_blocks.append(text)

#                 # ---- 表格 ----
#                 if shape.has_table:
#                     try:
#                         tbl = shape.table
#                         table_id = f"s{slide_num}_t{len(slide_tables)+1}"
#                         headers = [tbl.cell(0, c).text.strip() for c in range(len(tbl.columns))]
#                         rows = []
#                         for r in range(1, len(tbl.rows)):
#                             rows.append([
#                                 tbl.cell(r, c).text.strip()
#                                 for c in range(len(tbl.columns))
#                             ])
#                         coord_map = {}
#                         for ri in range(len(rows)):
#                             for ci in range(len(headers)):
#                                 coord_map[f"{ri},{ci}"] = f"slide{slide_num}"

#                         tb = TableBlock(
#                             meta=TableMeta(
#                                 table_id=table_id,
#                                 source_sheet=f"slide{slide_num}",
#                                 row_count=len(rows),
#                                 col_count=len(headers),
#                                 coord_map=coord_map,
#                             ),
#                             headers=headers,
#                             rows=rows,
#                         )
#                         slide_tables.append(tb)
#                         all_tables.append(tb)
#                     except Exception as exc:
#                         warnings.append(Warning(
#                             code="TABLE_EXTRACT_FAILED",
#                             message=f"Slide {slide_num} table extraction failed: {exc}",
#                             loc=slide_loc,
#                         ))

#                 # ---- 图片 → 只能提取引用，文字需 OCR ----
#                 try:
#                     if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
#                         image_count += 1
#                         img_id = f"s{slide_num}_img{image_count}"
#                         all_images.append(ImageBlock(
#                             image_id=img_id,
#                             caption="",
#                             loc=slide_loc,
#                         ))
#                         warnings.append(Warning(
#                             code="IMAGE_TEXT_NOT_EXTRACTED",
#                             message=(
#                                 f"Slide {slide_num}: image '{img_id}' detected. "
#                                 "Text inside images is NOT extracted. "
#                                 "Consider OCR or manual transcription."
#                             ),
#                             loc=slide_loc,
#                         ))
#                 except Exception:
#                     pass

#             # ---- 备注 ----
#             notes = ""
#             try:
#                 if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
#                     notes = slide.notes_slide.notes_text_frame.text.strip()
#             except Exception:
#                 pass

#             # 修正标题：如果占位符方式没拿到，用启发式
#             if not slide_title and text_blocks:
#                 slide_title = text_blocks.pop(0)

#             # 如果有 text_blocks 但没 bullets，把短文本也归入 bullets
#             if not bullets and text_blocks:
#                 for tb in text_blocks:
#                     bullets.append(tb)
#                 text_blocks = []

#             if slide_num == 1 and slide_title:
#                 title = slide_title

#             slides_data.append(SlideContent(
#                 slide_number=slide_num,
#                 title=slide_title,
#                 bullets=bullets,
#                 notes=notes,
#                 text_blocks=text_blocks,
#                 tables=slide_tables,
#                 loc=slide_loc,
#             ))

#         # 如果整个 PPT 文字非常少，发 warning
#         total_text = sum(
#             len(s.title) + sum(len(b) for b in s.bullets) + len(s.notes)
#             for s in slides_data
#         )
#         if total_text < 50 and len(slides_data) > 0:
#             warnings.append(Warning(
#                 code="LOW_TEXT_CONTENT",
#                 message=(
#                     f"Very little text extracted ({total_text} chars). "
#                     "This PPT may be image-heavy; consider OCR."
#                 ),
#                 loc=Loc(type=LocType.SLIDE.value, value=0),
#             ))

#         return ParsedDocument(
#             doc_id=doc_id,
#             source=self._source_dict(fpath),
#             title=title,
#             content_type="pptx",
#             slides=slides_data,
#             tables=all_tables,
#             images=all_images,
#             warnings=warnings,
#             errors=errors,
#             metadata={
#                 "total_slides": len(slides_data),
#                 "total_images": image_count,
#                 "total_text_chars": total_text,
#             },
#         )

#     def _source_dict(self, fpath: Path) -> dict:
#         return {
#             "file_name": fpath.name,
#             "file_type": "pptx",
#             "file_size_bytes": fpath.stat().st_size,
#             "file_path": str(fpath),
#         }
"""
PPT 解析器。
BOO-63 要求：
- 每页 Slide 输出：标题、bullets、notes
- 每条要点带 loc（至少 slide 编号）
- 不依赖大模型主观总结，先输出结构化要点
- 图片中文字使用 OCR 抽取，不发 warning
"""
from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from agent.parser.base import BaseParser
from agent.parser.schema import (
    Error, ImageBlock, Loc, LocType, ParsedDocument, SlideContent,
    TableBlock, TableMeta, Warning,
)
from agent.parser.utils.hash_utils import file_sha1
from agent.parser.utils.ocr_utils import run_ocr  # <--- 这里用你现有的函数

logger = logging.getLogger(__name__)


def _get_placeholder_title(slide) -> str:
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                text = ph.text_frame.text.strip() if ph.has_text_frame else ""
                if text:
                    return text
    except (ValueError, KeyError, AttributeError):
        pass
    return ""


class PptxParser(BaseParser):
    supported_extensions = [".pptx"]

    def __init__(self, kb_name: str = ""):
        self.kb_name = kb_name

    def parse(self, file_path: str, **kwargs) -> ParsedDocument:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        doc_id = file_sha1(file_path, prefix=self.kb_name)
        fpath = Path(file_path)
        warnings: list[Warning] = []
        errors: list[Error] = []

        try:
            prs = Presentation(file_path)
        except Exception as exc:
            return ParsedDocument(
                doc_id=doc_id,
                source=self._source_dict(fpath),
                content_type="pptx",
                errors=[Error(
                    error_code="PPTX_OPEN_FAILED",
                    message=str(exc),
                    loc=Loc(type=LocType.SLIDE.value, value=0),
                )],
            )

        slides_data: list[SlideContent] = []
        all_tables: list[TableBlock] = []
        all_images: list[ImageBlock] = []
        title = ""
        image_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_loc = Loc(type=LocType.SLIDE.value, value=slide_num)
            slide_title = _get_placeholder_title(slide)

            bullets: list[str] = []
            text_blocks: list[str] = []
            slide_tables: list[TableBlock] = []

            # ============================
            # 遍历形状：文本 + 表格 + 图片OCR
            # ============================
            for shape in slide.shapes:

                # ---------- 文本框 ----------
                if shape.has_text_frame:
                    is_title_shape = False
                    if shape.is_placeholder:
                        try:
                            if shape.placeholder_format.idx == 0:
                                is_title_shape = True
                        except (ValueError, AttributeError):
                            pass

                    if is_title_shape:
                        if not slide_title:
                            slide_title = shape.text_frame.text.strip()
                    else:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if not text:
                                continue
                            if para.level > 0 or text.startswith(("•", "-", "·", "▪", "●")):
                                bullets.append(text)
                            else:
                                if text != slide_title:
                                    text_blocks.append(text)

                # ---------- 表格 ----------
                if shape.has_table:
                    try:
                        tbl = shape.table
                        table_id = f"s{slide_num}_t{len(slide_tables)+1}"
                        headers = [tbl.cell(0, c).text.strip() for c in range(len(tbl.columns))]
                        rows = []
                        for r in range(1, len(tbl.rows)):
                            rows.append([
                                tbl.cell(r, c).text.strip()
                                for c in range(len(tbl.columns))
                            ])
                        coord_map = {}
                        for ri in range(len(rows)):
                            for ci in range(len(headers)):
                                coord_map[f"{ri},{ci}"] = f"slide{slide_num}"

                        tb = TableBlock(
                            meta=TableMeta(
                                table_id=table_id,
                                source_sheet=f"slide{slide_num}",
                                row_count=len(rows),
                                col_count=len(headers),
                                coord_map=coord_map,
                            ),
                            headers=headers,
                            rows=rows,
                        )
                        slide_tables.append(tb)
                        all_tables.append(tb)
                    except Exception as exc:
                        warnings.append(Warning(
                            code="TABLE_EXTRACT_FAILED",
                            message=f"Slide {slide_num} table extraction failed: {exc}",
                            loc=slide_loc,
                        ))

                # ---------- 图片 → OCR （无警告） ----------
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count += 1
                        img_id = f"s{slide_num}_img{image_count}"

                        # ========== OCR 核心调用 ==========
                        try:
                            img_bytes = shape.image.blob
                            ocr_text = run_ocr(img_bytes)
                            if ocr_text and ocr_text.strip():
                                text_blocks.append(f"[OCR] {ocr_text.strip()}")
                        except Exception as e:
                            logger.debug(f"OCR failed for {img_id}: {str(e)}")

                        # 只记录图片，不发警告
                        all_images.append(ImageBlock(
                            image_id=img_id,
                            caption="",
                            loc=slide_loc,
                        ))

                except Exception:
                    pass

            # ---------- 备注 ----------
            notes = ""
            try:
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                pass

            # ---------- 标题补全 ----------
            if not slide_title and text_blocks:
                slide_title = text_blocks.pop(0)

            if not bullets and text_blocks:
                for tb in text_blocks:
                    bullets.append(tb)
                text_blocks = []

            if slide_num == 1 and slide_title:
                title = slide_title

            slides_data.append(SlideContent(
                slide_number=slide_num,
                title=slide_title,
                bullets=bullets,
                notes=notes,
                text_blocks=text_blocks,
                tables=slide_tables,
                loc=slide_loc,
            ))

        # 低文本量警告（可保留，不影响）
        total_text = sum(
            len(s.title) + sum(len(b) for b in s.bullets) + len(s.notes)
            for s in slides_data
        )
        if total_text < 50 and len(slides_data) > 0:
            warnings.append(Warning(
                code="LOW_TEXT_CONTENT",
                message=f"Very little text extracted ({total_text} chars). OCR may help.",
                loc=Loc(type=LocType.SLIDE.value, value=0),
            ))

        return ParsedDocument(
            doc_id=doc_id,
            source=self._source_dict(fpath),
            title=title,
            content_type="pptx",
            slides=slides_data,
            tables=all_tables,
            images=all_images,
            warnings=warnings,  # <--- 无图片警告！
            errors=errors,
            metadata={
                "total_slides": len(slides_data),
                "total_images": image_count,
                "total_text_chars": total_text,
            },
        )

    def _source_dict(self, fpath: Path) -> dict:
        return {
            "file_name": fpath.name,
            "file_type": "pptx",
            "file_size_bytes": fpath.stat().st_size,
            "file_path": str(fpath),
        }