from __future__ import annotations

import argparse
import json
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, Iterable, List, Tuple

import pandas as pd
from docx import Document
from pptx import Presentation


SUPPORTED_SUFFIXES = {".txt", ".md", ".docx", ".pptx", ".xlsx", ".xls"}


def _clean_text(text: str) -> str:
    cleaned = text.replace("\u3000", " ")
    cleaned = re.sub(r"\r\n?", "\n", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    cleaned = re.sub(r"[ \t]{2,}", " ", cleaned)
    return cleaned.strip()


def _read_text_file(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return path.read_text(encoding="gbk", errors="ignore")


def _read_docx(path: Path) -> str:
    doc = Document(path)
    parts: List[str] = []

    for paragraph in doc.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n".join(parts)


def _read_pptx(path: Path) -> str:
    prs = Presentation(path)
    parts: List[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_texts: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text).strip()
                if text:
                    slide_texts.append(text)
        if slide_texts:
            parts.append(f"[Slide {slide_idx}]\n" + "\n".join(slide_texts))

    return "\n\n".join(parts)


def _read_excel(path: Path) -> str:
    xls = pd.ExcelFile(path)
    blocks: List[str] = []

    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name=sheet_name, dtype=str).fillna("")
        if df.empty:
            continue
        sheet_text = df.astype(str).apply(lambda row: " | ".join(row.tolist()), axis=1).tolist()
        blocks.append(f"[Sheet {sheet_name}]\n" + "\n".join(sheet_text))

    return "\n\n".join(blocks)


def _extract_text(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return _read_text_file(path)
    if suffix == ".docx":
        return _read_docx(path)
    if suffix == ".pptx":
        return _read_pptx(path)
    if suffix in {".xlsx", ".xls"}:
        return _read_excel(path)
    return ""


def _split_chunks(text: str, *, chunk_size: int, overlap: int) -> List[str]:
    normalized = _clean_text(text)
    if not normalized:
        return []

    paragraphs = [para.strip() for para in normalized.split("\n\n") if para.strip()]
    chunks: List[str] = []
    current = ""

    for para in paragraphs:
        candidate = f"{current}\n\n{para}".strip() if current else para
        if len(candidate) <= chunk_size:
            current = candidate
            continue

        if current:
            chunks.append(current)

        if len(para) <= chunk_size:
            current = para
        else:
            start = 0
            while start < len(para):
                end = min(start + chunk_size, len(para))
                piece = para[start:end].strip()
                if piece:
                    chunks.append(piece)
                if end >= len(para):
                    break
                start = max(end - overlap, start + 1)
            current = ""

    if current:
        chunks.append(current)

    return chunks


def _iter_files(source_dir: Path) -> Iterable[Path]:
    for path in source_dir.rglob("*"):
        if path.is_file() and path.suffix.lower() in SUPPORTED_SUFFIXES:
            yield path


def build_kb(
    source_dir: Path,
    output_file: Path,
    *,
    chunk_size: int,
    overlap: int,
) -> Tuple[int, int]:
    chunks: List[Dict[str, str]] = []
    file_count = 0

    for file_path in _iter_files(source_dir):
        text = _extract_text(file_path)
        file_chunks = _split_chunks(text, chunk_size=chunk_size, overlap=overlap)
        if not file_chunks:
            continue

        file_count += 1
        relative_source = file_path.relative_to(source_dir).as_posix()

        for idx, chunk in enumerate(file_chunks, start=1):
            chunks.append(
                {
                    "id": f"{relative_source}#{idx}",
                    "source": relative_source,
                    "title": f"{file_path.stem}-片段{idx}",
                    "content": chunk,
                }
            )

    payload = {
        "metadata": {
            "source_dir": str(source_dir),
            "built_at": datetime.now().isoformat(timespec="seconds"),
            "chunk_size": chunk_size,
            "overlap": overlap,
            "file_count": file_count,
            "chunk_count": len(chunks),
        },
        "chunks": chunks,
    }

    output_file.parent.mkdir(parents=True, exist_ok=True)
    output_file.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")

    # Persist to chromadb
    try:
        import chromadb
        from chromadb.config import Settings
        from chromadb.utils import embedding_functions

        db_path = output_file.parent / "chroma_db"
        client = chromadb.PersistentClient(path=str(db_path))

        # We'll stick to a simple sentence-transformer model in chroma instead of loading our custom one,
        # or we could make our custom embedder to be used with chroma.
        # But this suffices for an ingest demo that chroma does the embedding locally
        # You specified 'jinaai/jina-embeddings-v5-text-nano-retrieval' in retriever.py, let's use it correctly.
        class JinaEmbeddingFunction(embedding_functions.EmbeddingFunction):
            def __call__(self, input: chromadb.Documents) -> chromadb.Embeddings:
                import torch
                from sentence_transformers import SentenceTransformer
                device = "cuda" if torch.cuda.is_available() else "cpu"
                model = SentenceTransformer(
                    "jinaai/jina-embeddings-v5-text-nano-retrieval",
                    trust_remote_code=True,
                    device=device,
                )
                embeddings = model.encode(input, convert_to_numpy=True, show_progress_bar=False, batch_size=32)
                return embeddings.tolist()
                
        emb_fn = JinaEmbeddingFunction()

        collection = client.get_or_create_collection(
            name="reimbursement_kb",
            embedding_function=emb_fn
        )
        
        # Insert them into Chroma
        if chunks:
            ids = [c["id"] for c in chunks]
            documents = [(c["title"] + "\\n" + c["content"]).strip() for c in chunks]
            metadatas = [{"source": c["source"], "title": c["title"], "content": c["content"]} for c in chunks]
            
            # Delete old content first to avoid duplicates across runs
            collection.delete(ids=ids)
            collection.upsert(
                documents=documents,
                metadatas=metadatas,
                ids=ids
            )
    except ImportError:
        pass

    return file_count, len(chunks)


def main() -> None:
    parser = argparse.ArgumentParser(description="构建报销知识库索引")
    parser.add_argument("--source", default="docs/reimbursement", help="资料目录")
    parser.add_argument("--output", default="data/kb/reimbursement_kb.json", help="索引输出路径")
    parser.add_argument("--chunk-size", type=int, default=700, help="单片段最大字符数")
    parser.add_argument("--overlap", type=int, default=100, help="分片重叠字符数")

    args = parser.parse_args()

    source_dir = Path(args.source).resolve()
    output_file = Path(args.output).resolve()

    if not source_dir.exists() or not source_dir.is_dir():
        raise FileNotFoundError(f"资料目录不存在: {source_dir}")

    file_count, chunk_count = build_kb(
        source_dir=source_dir,
        output_file=output_file,
        chunk_size=max(args.chunk_size, 200),
        overlap=max(min(args.overlap, args.chunk_size - 1), 0),
    )

    print(
        json.dumps(
            {
                "ok": True,
                "source": str(source_dir),
                "output": str(output_file),
                "file_count": file_count,
                "chunk_count": chunk_count,
            },
            ensure_ascii=False,
            indent=2,
        )
    )


if __name__ == "__main__":
    main()
